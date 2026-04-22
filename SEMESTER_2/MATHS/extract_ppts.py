from pptx import Presentation
import os, glob

ppt_dir = "Lecture Notes for CSDA Students-20260418"
ppts = sorted(glob.glob(os.path.join(ppt_dir, "*.pptx")))

for ppt_path in ppts:
    print(f"\n{'='*80}")
    print(f"FILE: {os.path.basename(ppt_path)}")
    print(f"{'='*80}")
    try:
        prs = Presentation(ppt_path)
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            texts.append(t)
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = " | ".join(cell.text.strip() for cell in row.cells)
                        if row_text.strip().replace("|","").strip():
                            texts.append(f"TABLE: {row_text}")
            if texts:
                print(f"\n--- Slide {i} ---")
                for t in texts:
                    print(t)
    except Exception as e:
        print(f"ERROR: {e}")
